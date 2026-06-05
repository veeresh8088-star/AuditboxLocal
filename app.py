# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import requests
import pdfplumber
import time, json, hashlib, uuid, threading
import logging
from datetime import datetime, timedelta

logging.basicConfig(level=logging.INFO, format='%(message)s')

from docx import Document
from database import engine, db_label, AuditFinding, ChatMessage, AuditCheckpoint, SessionLocal
from auth import render_login_gate
import scoping_engine
import easyocr
import numpy as np
import PIL.Image

# Thread-safe storage for background analysis results and active runs
@st.cache_resource
def _get_bg_store():
    return {
        "results": {},
        "running": set(),
        "progress": {},
        "lock": threading.Lock()
    }

_bg_store = _get_bg_store()
_bg_results = _bg_store["results"]
_bg_running = _bg_store["running"]
_bg_lock = _bg_store["lock"]

@st.cache_resource
def get_ocr_reader():
    # Only load English models into memory when needed
    return easyocr.Reader(['en'], gpu=False)

st.set_page_config(page_title="AICyberAuditBox", page_icon="🛡️", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');
html, body, [class*="css"] { font-family: 'Inter', sans-serif !important; }

/* ── Sidebar Primary Buttons ── */
section[data-testid="stSidebar"] button[data-testid="stBaseButton-primary"] {
    background: #3b82f6 !important;
    color: white !important;
    border: none !important;
    border-radius: 8px !important;
    font-weight: 600 !important;
    transition: 0.2s !important;
}
section[data-testid="stSidebar"] button[data-testid="stBaseButton-primary"]:hover {
    background: #2563eb !important;
    box-shadow: 0 4px 12px rgba(59,130,246,0.4) !important;
}

/* ── ChatGPT Style Recents Sidebar CSS ── */
.chatgpt-sidebar-list {
    display: flex;
    flex-direction: column;
    gap: 4px;
    padding: 0px !important;
}

.chatgpt-row {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 8px 10px;
    border-radius: 8px;
    transition: background-color 0.2s, border-color 0.2s;
    position: relative;
    margin-bottom: 2px;
}

.chatgpt-row-inactive {
    background-color: transparent;
    border: 1px solid transparent;
}

.chatgpt-row-inactive:hover {
    background-color: rgba(128, 128, 128, 0.08) !important; /* Theme-adaptive hover overlay */
}

.chatgpt-row-active {
    background-color: rgba(128, 128, 128, 0.15) !important; /* Theme-adaptive active background */
    border: 1px solid rgba(128, 128, 128, 0.25) !important; /* Theme-adaptive active border */
}

.chatgpt-row-left {
    display: flex;
    align-items: center;
    gap: 8px;
    min-width: 0;
    flex: 1;
    text-decoration: none !important;
}

.chatgpt-row-icon {
    font-size: 13px;
    color: var(--secondary-text-color) !important; /* Theme-adaptive secondary text */
    opacity: 0.85 !important;
    flex-shrink: 0;
    display: flex;
    align-items: center;
}

.chatgpt-row-title {
    font-size: 13px;
    font-weight: 500;
    color: var(--text-color) !important; /* Theme-adaptive main text */
    opacity: 0.8 !important; /* Muted opacity for inactive chats */
    white-space: nowrap;
    overflow: hidden;
    text-overflow: ellipsis;
    text-decoration: none !important;
}

.chatgpt-row-active .chatgpt-row-title {
    color: var(--text-color) !important;
    opacity: 1 !important; /* Full contrast for active title */
    font-weight: 600 !important;
}

.chatgpt-row-active .chatgpt-row-icon {
    opacity: 1 !important;
}

.chatgpt-row-delete {
    display: flex;
    align-items: center;
    justify-content: center;
    opacity: 0;
    transition: opacity 0.15s;
    z-index: 10;
    margin-left: 6px;
}

.chatgpt-row:hover .chatgpt-row-delete {
    opacity: 1;
}

.chatgpt-row-delete-link {
    color: var(--secondary-text-color) !important;
    text-decoration: none !important;
    font-size: 12px;
    width: 20px;
    height: 20px;
    border-radius: 4px;
    display: flex;
    align-items: center;
    justify-content: center;
    transition: background-color 0.15s, color 0.15s;
}

.chatgpt-row-delete-link:hover {
    background-color: rgba(239, 68, 68, 0.15) !important;
    color: #ef4444 !important;
}

/* ── Main UI styles ── */
.main-header { background:#1e293b;
    padding:28px 32px; border-radius:16px; margin-bottom:24px;
    border:1px solid rgba(59,130,246,.2); }
.stat-card { background:#1e293b; border:1px solid #334155; border-radius:12px;
    padding:20px; text-align:center; }
.stat-num { font-size:2rem; font-weight:700; }
.badge-critical { background:#1a0a0a; border:1px solid #ef4444; border-left:5px solid #ef4444; border-radius:8px; padding:16px; margin:8px 0; color:#f8fafc; }
.badge-high     { background:#1a0d08; border:1px solid #f97316; border-left:5px solid #f97316; border-radius:8px; padding:16px; margin:8px 0; color:#f8fafc; }
.badge-medium   { background:#1a1600; border:1px solid #eab308; border-left:5px solid #eab308; border-radius:8px; padding:16px; margin:8px 0; color:#f8fafc; }
.badge-low      { background:#051a0d; border:1px solid #22c55e; border-left:5px solid #22c55e; border-radius:8px; padding:16px; margin:8px 0; color:#f8fafc; }
.chat-bubble-user { background:#1e3a5f; border-radius:16px 16px 4px 16px; padding:12px 16px;
    margin:4px 0; max-width:80%; color:#e2e8f0; text-align:left; }
.chat-bubble-bot  { background:#1e293b; border:1px solid #334155; border-radius:16px 16px 16px 4px;
    padding:12px 16px; margin:4px 0; max-width:80%; color:#e2e8f0; text-align:left; }
.uc-card { background:#1e293b; border:1px solid #334155; border-radius:10px;
    padding:14px 18px; margin:8px 0; cursor:pointer; transition:.2s; }
.uc-card:hover { border-color:#3b82f6; transform:translateX(4px); }
.stage-done   { color:#22c55e; border-left:3px solid #22c55e; padding:6px 0 6px 14px; margin:4px 0; font-weight:600; }
.stage-active { color:#3b82f6; border-left:3px solid #3b82f6; padding:6px 0 6px 14px; margin:4px 0; font-weight:600; }
.stage-idle   { color:#475569; border-left:3px solid #334155; padding:6px 0 6px 14px; margin:4px 0; }
div[data-testid="stDecoration"] { display:none; }
.inline-spinner { border: 2px solid rgba(59, 130, 246, 0.1); border-top: 2px solid #3b82f6; border-radius: 50%; width: 16px; height: 16px; animation: spin_inline 1s linear infinite; display: inline-block; vertical-align: middle; }
@keyframes spin_inline { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }
</style>
""", unsafe_allow_html=True)

# ── AUTHENTICATION GATE ───────────────────────────────────────────────────────
render_login_gate()

# ── USE CASES ─────────────────────────────────────────────────────────────────
from controls_data import USE_CASES, DEMO_FINDINGS, GAP_RESOLUTION, SCOPE_KEYWORDS

def save_findings(uc, findings):
    db = SessionLocal()
    db.query(AuditFinding).filter(AuditFinding.use_case_sl == uc["sl"]).delete()
    uc_name = uc.get("use_case", uc.get("label", "Comprehensive Enterprise Audit"))
    for f in findings:
        db.add(AuditFinding(
            use_case_sl=uc["sl"],
            use_case_name=uc_name[:290],
            control_id=f.get("control_id", ""),
            relevance_score=f.get("relevance_score", 0),
            evidence_found=f.get("evidence_found", ""),
            evidence_snippet=f.get("evidence_snippet", ""),
            severity=f.get("severity", ""),
            control=f.get("control", ""),
            finding=f.get("finding", ""),
            recommendation=f.get("recommendation", ""),
            reasoning=f.get("reasoning", ""),
            status=f.get("status", "Open"),
            comment=f.get("comment", ""),
            source_files=f.get("source_files", ""),
        ))
    db.commit(); db.close()

def get_all_findings():
    db = SessionLocal()
    rows = db.query(AuditFinding).order_by(AuditFinding.created_at.desc()).all()
    db.close(); return rows

def save_chat_message(session_id, session_title, role, content):
    db = SessionLocal()
    db.query(ChatMessage).filter(ChatMessage.session_id == session_id).update({ChatMessage.session_title: session_title})
    db.add(ChatMessage(session_id=session_id, session_title=session_title, role=role, content=content))
    db.commit()
    db.close()

def update_latest_assistant_message(session_id, content):
    db = SessionLocal()
    latest = db.query(ChatMessage).filter(
        ChatMessage.session_id == session_id,
        ChatMessage.role == "assistant"
    ).order_by(ChatMessage.created_at.desc()).first()
    if latest:
        latest.content = content
        db.commit()
    db.close()

def get_chat_history(session_id):
    db = SessionLocal()
    msgs = db.query(ChatMessage).filter(ChatMessage.session_id == session_id).order_by(ChatMessage.created_at.asc()).all()
    db.close()
    return [{"role": m.role, "content": m.content} for m in msgs]

def get_chat_title(session_id):
    db = SessionLocal()
    msg = db.query(ChatMessage).filter(ChatMessage.session_id == session_id).order_by(ChatMessage.created_at.asc()).first()
    db.close()
    return msg.session_title if msg else None

def get_all_chat_sessions():
    db = SessionLocal()
    rows = db.query(
        ChatMessage.session_id,
        ChatMessage.session_title,
        ChatMessage.created_at
    ).order_by(ChatMessage.created_at.desc()).all()
    db.close()
    seen = set()
    sessions = []
    for r in rows:
        if r.session_id not in seen:
            seen.add(r.session_id)
            sessions.append({
                "session_id": r.session_id,
                "session_title": r.session_title,
                "created_at": r.created_at
            })
            if len(sessions) == 10:
                break
    return sessions

def clear_chat_session(session_id):
    db = SessionLocal()
    db.query(ChatMessage).filter(ChatMessage.session_id == session_id).delete()
    db.commit()
    db.close()

# ── CHECKPOINT HELPERS ────────────────────────────────────────────────────────
def _checkpoint_create(session_id, bg_key, ai_model, selected_sls, file_names, context_str, total_controls, batch_size):
    """Create a fresh in-progress checkpoint row when an audit starts."""
    db = SessionLocal()
    try:
        # Remove any stale checkpoint for this session
        db.query(AuditCheckpoint).filter(AuditCheckpoint.session_id == session_id).delete()
        chk = AuditCheckpoint(
            session_id=session_id,
            bg_key=bg_key,
            ai_model=ai_model,
            selected_sls_json=json.dumps(list(selected_sls)),
            file_names_json=json.dumps(file_names),
            context_text=context_str,
            total_controls=total_controls,
            completed_batches=0,
            batch_size=batch_size,
            partial_results_json="[]",
            status="in_progress",
        )
        db.add(chk)
        db.commit()
        return chk.id
    except Exception as e:
        print(f"[checkpoint] Failed to create checkpoint: {e}")
        return None
    finally:
        db.close()

def _checkpoint_update(session_id, completed_batches, all_results_so_far):
    """Persist partial results after each batch completes."""
    db = SessionLocal()
    try:
        chk = db.query(AuditCheckpoint).filter(
            AuditCheckpoint.session_id == session_id,
            AuditCheckpoint.status == "in_progress"
        ).first()
        if chk:
            chk.completed_batches = completed_batches
            chk.partial_results_json = json.dumps(all_results_so_far)
            chk.updated_at = datetime.utcnow()
            db.commit()
    except Exception as e:
        print(f"[checkpoint] Failed to update checkpoint: {e}")
    finally:
        db.close()

def _checkpoint_finish(session_id, status="completed"):
    """Mark the checkpoint as completed or failed."""
    db = SessionLocal()
    try:
        chk = db.query(AuditCheckpoint).filter(
            AuditCheckpoint.session_id == session_id,
            AuditCheckpoint.status == "in_progress"
        ).first()
        if chk:
            chk.status = status
            chk.updated_at = datetime.utcnow()
            db.commit()
    except Exception as e:
        print(f"[checkpoint] Failed to finish checkpoint: {e}")
    finally:
        db.close()

def get_resumable_checkpoint(session_id):
    """Return an in-progress checkpoint for this session, or None."""
    db = SessionLocal()
    try:
        return db.query(AuditCheckpoint).filter(
            AuditCheckpoint.session_id == session_id,
            AuditCheckpoint.status == "in_progress"
        ).order_by(AuditCheckpoint.created_at.desc()).first()
    finally:
        db.close()

def extract_text(f):
    name_lower = f.name.lower()

    # ── ZIP / Folder Upload ─────────────────────────────────────────────────
    if name_lower.endswith(".zip"):
        import zipfile, io as _io
        SUPPORTED_EXTS = (
            ".pdf", ".docx", ".doc", ".xlsx", ".xls",
            ".csv", ".pptx", ".ppt", ".txt",
            ".png", ".jpg", ".jpeg"
        )
        combined_texts = []
        try:
            with zipfile.ZipFile(_io.BytesIO(f.read())) as zf:
                # Sort so parent dirs are processed before children
                entries = sorted(zf.namelist())
                for entry in entries:
                    # Skip directories, hidden files, macOS metadata
                    if entry.endswith("/") or "__MACOSX" in entry or entry.startswith("."):
                        continue
                    entry_lower = entry.lower()
                    if not any(entry_lower.endswith(ext) for ext in SUPPORTED_EXTS):
                        continue   # skip unsupported file types silently
                    try:
                        with zf.open(entry) as inner_file:
                            inner_bytes = inner_file.read()
                        inner_f = _io.BytesIO(inner_bytes)
                        # Give it a .name attribute so extract_text sub-calls work
                        inner_f.name = entry.split("/")[-1]  # use basename
                        inner_text = extract_text(inner_f)
                        if inner_text and not inner_text.startswith("[Error"):
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n{inner_text}")
                        elif inner_text.startswith("[Error"):
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n{inner_text}")
                    except Exception as ie:
                        combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n[Error reading {entry}: {ie}]")
            if combined_texts:
                return "\n\n".join(combined_texts)
            return "[ZIP file appears empty or contains no supported document types.]"
        except zipfile.BadZipFile:
            return f"[Error: {f.name} is not a valid ZIP file.]"
        except Exception as e:
            return f"[Error extracting ZIP {f.name}: {e}]"

    # ── Image files (PNG / JPG / JPEG) ──────────────────────────────────────
    if name_lower.endswith((".png", ".jpg", ".jpeg")):
        try:
            reader = get_ocr_reader()
            img = PIL.Image.open(f)
            img_np = np.array(img)
            res = reader.readtext(img_np, detail=0)
            return " ".join(res)
        except Exception as e:
            return f"[Error parsing image file {f.name}: {e}]"

    elif name_lower.endswith(".pdf"):
        with pdfplumber.open(f) as pdf:
            pages_text = []
            for p in pdf.pages:
                text = p.extract_text() or ""
                # If native text is extremely short (likely a scanned PDF image)
                if len(text.strip()) < 50:
                    try:
                        reader = get_ocr_reader()
                        img = p.to_image(resolution=200).original
                        img_np = np.array(img)
                        res = reader.readtext(img_np, detail=0)
                        text += " " + " ".join(res)
                    except:
                        pass
                pages_text.append(text)
            return "\n".join(pages_text)
    elif name_lower.endswith((".xlsx", ".xls")):
        try:
            excel_data = pd.read_excel(f, sheet_name=None)
            sheets_text = []
            for sheet_name, df in excel_data.items():
                sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + df.to_string(index=False))
            return "\n\n".join(sheets_text)
        except Exception as e:
            return f"[Error parsing Excel file {f.name}: {e}]"
    elif name_lower.endswith(".csv"):
        try:
            df = pd.read_csv(f)
            return df.to_string(index=False)
        except Exception as e:
            return f"[Error parsing CSV file {f.name}: {e}]"
    elif name_lower.endswith((".pptx", ".ppt")):
        try:
            from pptx import Presentation
            prs = Presentation(f)
            text_runs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_runs.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                text_runs.append(" | ".join(row_text))
            return "\n".join(text_runs)
        except Exception as e:
            return f"[Error parsing PowerPoint file {f.name}: {e}]"
    elif name_lower.endswith(".txt"):
        try:
            return f.read().decode("utf-8", errors="ignore")
        except Exception as e:
            return f"[Error parsing text file {f.name}: {e}]"
    else:
        try:
            doc = Document(f)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            return f"[Error parsing Word file {f.name}: {e}]"

def scan_file_security(uploaded_file):
    bytes_data = uploaded_file.getvalue()
    uploaded_file.seek(0)
    if bytes_data.startswith(b'MZ'):
        return False, "Executable payload disguised as document (MZ signature detected)."
    file_hash = hashlib.sha256(bytes_data).hexdigest()
    blacklist = ["5e884898da28047151d0e56f8dc6292773603d0d6aabbdd62a11ef721d1542d8"]
    if file_hash in blacklist:
        return False, f"Known malicious hash signature identified ({file_hash[:8]}...)."
    return True, "Clean"

def _resolve_ollama_model(model_choice):
    """Map UI model name to Ollama model identifier."""
    if "7B" in model_choice:
        return "qwen2.5:7b"
    elif "3B" in model_choice:
        return "qwen2.5:3b"
    elif "Gemma 2 (9B)" in model_choice:
        return "gemma2:9b"
    return "llama3.1:latest"


def _build_controls_for_audit(selected_sls):
    """Gather control metadata for the selected sl numbers."""
    controls = []
    for uc in USE_CASES:
        if uc["sl"] in selected_sls:
            controls.append({
                "control": uc["use_case"],
                "label": uc["label"],
                "expected": uc["expected"],
                "prompt_hint": uc["prompt_hint"],
                "severity": uc.get("severity", "MEDIUM"),
                "standard": uc.get("standard", ""),
            })
    return controls


def _audit_batch(context, controls_batch, file_names_list, ollama_model, timeout=240):
    """Send a single batch of controls to the LLM for structured ISO 27001 audit.

    Returns a list of result dicts, one per control, each containing:
      control_id, control, relevance_score, evidence_found, evidence_snippet,
      status, severity, finding, recommendation, reasoning, source_files
    """
    scanned_files_str = ", ".join(file_names_list) if file_names_list else "None"

    controls_desc = []
    for i, c in enumerate(controls_batch, 1):
        controls_desc.append(
            f"{i}. Control ID: {c['control']}\n"
            f"   Description: {c['label']}\n"
            f"   Expected evidence: {c['expected']}\n"
            f"   Audit instruction: {c['prompt_hint']}"
        )
    controls_text = "\n".join(controls_desc)

    prompt = f"""You are a Senior ISO 27001 Lead Auditor with 15+ years of experience.
Perform an accurate compliance audit with minimal false positives and false negatives.

AUDIT PRINCIPLES:
1. Think like a human auditor — do NOT audit controls blindly.
2. Use semantic understanding — do NOT require exact ISO wording.
3. Search for direct, supporting, related, AND implied evidence.
4. If evidence exists, do not mark Non-Compliant.
5. Mark Partial Compliance when SOME evidence exists.
6. Mark Out Of Scope when the control is NOT relevant to this document.
7. Always explain your reasoning.

SCOPING & RELEVANCE:
- Calculate a relevance_score (0-100) for each control against this document.
- Score >= 80: Audit fully.
- Score 60-79: Audit only if evidence exists.
- Score < 60: Mark Out Of Scope — do not force a finding.

EVIDENCE EVALUATION:
- Strong Evidence (direct, explicit): status = "Compliant"
- Some Evidence (partial, implied): status = "Partially Compliant"
- No Evidence (nothing found): status = "Non-Compliant"
- Control not relevant to document: status = "Out Of Scope"

SEVERITY (for Non-Compliant and Partially Compliant only):
- P1 Critical: Major compliance gap
- P2 High: Significant weakness
- P3 Medium: Partial implementation
- P4 Low: Minor documentation issue

FALSE POSITIVE PREVENTION:
- Never create findings solely because keywords are missing.
- Use semantic reasoning — if a procedure exists under a different name, it still counts.

FALSE NEGATIVE PREVENTION:
- Search for: Direct Evidence, Supporting Evidence, Related Evidence, Implied Evidence.

EVIDENCE TEXT (from: {scanned_files_str}):
\"\"\"
{context[:14000]}
\"\"\"

CONTROLS TO AUDIT:
{controls_text}

INSTRUCTIONS — follow EXACTLY:
For EACH control, produce a JSON object with these fields:
  - "control_id": the control ID string (e.g. "ISO-1215 Access Control")
  - "control": human-readable control name (e.g. "Access Control")
  - "relevance_score": integer 0-100
  - "evidence_found": one of "Strong Evidence", "Some Evidence", "No Evidence", "Not Relevant"
  - "evidence_snippet": a short direct quote or description of the evidence found (max 2 sentences). Empty string if none.
  - "status": one of "Compliant", "Partially Compliant", "Non-Compliant", "Out Of Scope"
  - "severity": "P1 Critical", "P2 High", "P3 Medium", or "P4 Low" (use "N/A" for Compliant or Out Of Scope)
  - "finding": 2-4 sentence explanation of what was found or what is missing.
  - "recommendation": specific actionable steps. Empty string if Compliant or Out Of Scope.
  - "reasoning": 1-3 sentences explaining your audit reasoning and why you assigned this status.
  - "source_files": "{scanned_files_str}"

Return ONLY valid JSON — no markdown, no extra text:
{{
  "results": [
    {{
      "control_id": "...",
      "control": "...",
      "relevance_score": 85,
      "evidence_found": "Strong Evidence",
      "evidence_snippet": "...",
      "status": "Compliant",
      "severity": "N/A",
      "finding": "...",
      "recommendation": "",
      "reasoning": "...",
      "source_files": "..."
    }}
  ]
}}
"""
    try:
        r = requests.post(
            "http://127.0.0.1:11434/api/generate",
            json={"model": ollama_model, "prompt": prompt, "stream": False, "format": "json"},
            timeout=timeout,
        )
        if r.status_code == 200:
            res = r.json().get("response", "{}")
            data = json.loads(res)
            results = data.get("results", [])
            if isinstance(results, list):
                return results
    except Exception as e:
        print(f"[_audit_batch] LLM call error: {e}")
    return None


def generate_ollama_findings(context, file_names_list, selected_sls, model_choice, bg_key=None, batch_size=10, checkpoint_session_id=None):
    """Audit controls in batches and return (resolved_list, findings).

    Uses the ISO 27001 Lead Auditor logic:
      Compliant       -> resolved_list
      Out Of Scope    -> skipped (not a finding)
      Partially Compliant / Non-Compliant -> findings

    If checkpoint_session_id is set, partial results are saved to ShaktiDB
    after every batch so the audit can be resumed after a crash.
    """
    ollama_model = _resolve_ollama_model(model_choice)
    controls = _build_controls_for_audit(selected_sls)

    if not controls:
        return [], []

    scanned_files_str = ", ".join(file_names_list) if file_names_list else "None"
    all_results = []
    total = len(controls)

    # Split into batches
    batches = [controls[i:i + batch_size] for i in range(0, total, batch_size)]

    import time
    overall_start_time = time.time()
    print(f"\n[{time.strftime('%H:%M:%S')}] [INFO] Starting ISO 27001 Audit for {total} controls in {len(batches)} batches...")

    # Valid severity values from the new prompt
    VALID_SEVERITIES = ("P1 Critical", "P2 High", "P3 Medium", "P4 Low", "N/A")
    # Map old uppercase values to new P-notation (for LLM fallback)
    SEV_UPGRADE_MAP = {"CRITICAL": "P1 Critical", "HIGH": "P2 High", "MEDIUM": "P3 Medium", "LOW": "P4 Low"}

    for batch_idx, batch in enumerate(batches):
        start_n = batch_idx * batch_size + 1
        end_n = min(start_n + batch_size - 1, total)

        batch_start_time = time.time()
        print(f"[{time.strftime('%H:%M:%S')}]   -> Running Batch {batch_idx + 1}/{len(batches)} (controls {start_n} to {end_n})...")

        if bg_key:
            with _bg_lock:
                _bg_store["progress"][bg_key] = f"⚡ Auditing controls {start_n}–{end_n} of {total}..."

        batch_results = _audit_batch(context, batch, file_names_list, ollama_model)

        if batch_results is None:
            # Fallback: treat all controls in this batch as Non-Compliant with generic findings
            for c in batch:
                all_results.append({
                    "control_id": c["control"],
                    "control": c["label"],
                    "relevance_score": 50,
                    "evidence_found": "No Evidence",
                    "evidence_snippet": "",
                    "status": "Non-Compliant",
                    "severity": SEV_UPGRADE_MAP.get(c.get("severity", "MEDIUM").upper(), "P3 Medium"),
                    "finding": f"No documented evidence found for {c['control']} ({c['label']}). LLM batch call failed.",
                    "recommendation": f"Establish, document, and implement procedures to satisfy {c['control']} ({c['label']}).",
                    "reasoning": "Fallback result — LLM call failed for this batch.",
                    "source_files": f"Checked in: {scanned_files_str}",
                })
        else:
            # Build lookup by control_id (primary) or control name (fallback)
            returned_by_id   = {r.get("control_id", ""): r for r in batch_results}
            returned_by_name = {r.get("control", ""):    r for r in batch_results}

            for c in batch:
                result = returned_by_id.get(c["control"]) or returned_by_name.get(c["label"])

                if result:
                    # Normalize status
                    raw_status = result.get("status", "Non-Compliant")
                    if raw_status not in ("Compliant", "Partially Compliant", "Non-Compliant", "Out Of Scope"):
                        # Map old Resolved/Unresolved to new schema
                        if raw_status == "Resolved":
                            raw_status = "Compliant"
                        elif raw_status == "Unresolved":
                            raw_status = "Non-Compliant"
                        else:
                            raw_status = "Non-Compliant"
                    result["status"] = raw_status

                    # Normalize severity
                    raw_sev = result.get("severity", "P3 Medium")
                    if raw_sev.upper() in SEV_UPGRADE_MAP:
                        raw_sev = SEV_UPGRADE_MAP[raw_sev.upper()]
                    if raw_sev not in VALID_SEVERITIES:
                        raw_sev = SEV_UPGRADE_MAP.get(c.get("severity", "MEDIUM").upper(), "P3 Medium")
                    result["severity"] = raw_sev

                    # Ensure all new fields are present
                    result.setdefault("control_id",       c["control"])
                    result.setdefault("control",          result.get("control", c["label"]))
                    result.setdefault("relevance_score",  50)
                    result.setdefault("evidence_found",   "No Evidence")
                    result.setdefault("evidence_snippet", "")
                    result.setdefault("reasoning",        "")
                    result.setdefault("source_files",     scanned_files_str)
                    all_results.append(result)
                else:
                    all_results.append({
                        "control_id": c["control"],
                        "control": c["label"],
                        "relevance_score": 50,
                        "evidence_found": "No Evidence",
                        "evidence_snippet": "",
                        "status": "Non-Compliant",
                        "severity": SEV_UPGRADE_MAP.get(c.get("severity", "MEDIUM").upper(), "P3 Medium"),
                        "finding": f"No documented evidence found for {c['control']} ({c['label']}).",
                        "recommendation": f"Establish, document, and implement procedures to satisfy {c['control']} ({c['label']}).",
                        "reasoning": "Control not returned by LLM. Defaulting to Non-Compliant.",
                        "source_files": scanned_files_str,
                    })

        batch_elapsed = time.time() - batch_start_time
        print(f"[{time.strftime('%H:%M:%S')}]   [SUCCESS] Batch {batch_idx + 1} completed in {batch_elapsed:.2f}s")

        # ── Persist checkpoint after every batch ───────────────────────────────
        if checkpoint_session_id:
            _checkpoint_update(checkpoint_session_id, batch_idx + 1, all_results)

    overall_elapsed = time.time() - overall_start_time
    print(f"[{time.strftime('%H:%M:%S')}] [SUCCESS] ISO 27001 Audit complete! Total time: {overall_elapsed:.2f} seconds.")

    # Compliant  -> resolved list
    # Out Of Scope -> excluded (not a finding, not resolved)
    # Partially Compliant / Non-Compliant -> findings
    resolved_list = [r["control_id"] for r in all_results if r.get("status") == "Compliant"]
    findings = [
        r for r in all_results
        if r.get("status") in ("Partially Compliant", "Non-Compliant")
    ]

    oos_count = sum(1 for r in all_results if r.get("status") == "Out Of Scope")
    if oos_count:
        print(f"   [INFO] {oos_count} control(s) marked Out Of Scope — excluded from findings.")

    return resolved_list, findings

def ai_chat_stream(system_ctx, user_msg, model_choice):
    enhanced_sys = f"You are a Senior Cybersecurity Auditor with expertise in ISO 27001, NIST, and SOC 2. {system_ctx}"
    prompt = f"{enhanced_sys}\n\nUser: {user_msg}\n\nAI Auditor:"
    if "Escalation" in model_choice:
        ollama_model = "qwen2.5:7b"
    else:
        ollama_model = _resolve_ollama_model(model_choice)
    try:
        r = requests.post("http://127.0.0.1:11434/api/generate",
            json={"model": ollama_model, "prompt": prompt, "stream": True}, stream=True, timeout=90)
        if r.status_code != 200:
            try:
                err = r.json().get("error", r.text)
            except:
                err = r.text
            yield f"⚠️ Ollama Error: {err}. Please make sure you have downloaded the model using pull_models.bat!"
            return
        for line in r.iter_lines():
            if line:
                chunk = json.loads(line)
                yield chunk.get("response", "")
    except Exception as e:
        yield f"⚠️ Offline Engine not responding: {e}"

def _run_ollama_bg(bg_key, files_data, selected_sls_copy, ai_model, session_id=None):
    import io
    print(f"[_run_ollama_bg] Starting thread for key {bg_key} with model {ai_model}...")
    _sid = session_id or bg_key   # use session_id for checkpoint keying
    try:
        with _bg_lock:
            _bg_store["progress"][bg_key] = "🔍 Scanning file security..."
        ctx = ""
        file_names_list = []
        for f_data in files_data:
            name = f_data["name"]
            file_bytes = f_data["bytes"]
            f_like = io.BytesIO(file_bytes)
            f_like.name = name
            is_clean, reason = scan_file_security(f_like)
            if not is_clean:
                print(f"[_run_ollama_bg] Security alert! Malware scan failed for file {name}: {reason}")
                with _bg_lock:
                    _bg_results[bg_key] = {"error": f"🚨 SECURITY ALERT: '{name}' BLOCKED! {reason}"}
                    _bg_store["progress"].pop(bg_key, None)
                _checkpoint_finish(_sid, "failed")
                return
            text = extract_text(f_like)
            ctx += f"--- FILE: {name} ---\n{text}\n\n"
            file_names_list.append(name)
        context_str = ctx.strip()

        # ── Create checkpoint so we can resume if the process crashes ─────────
        from controls_data import USE_CASES as _UC
        _total_ctrl_count = len([u for u in _UC if u["sl"] in selected_sls_copy])
        _batch_sz = 10
        _checkpoint_create(
            _sid, bg_key, ai_model,
            selected_sls_copy, file_names_list, context_str,
            _total_ctrl_count, _batch_sz
        )
        print(f"[checkpoint] Created checkpoint for session {_sid}")

        if ai_model == "Escalation Mode (Qwen 3B -> 7B) - High Accuracy/Reasoning":
            # Pass 1: fast scan with 3B
            with _bg_lock:
                _bg_store["progress"][bg_key] = "⚡ Pass 1/2 — Qwen 2.5 (3B) fast-pass scan..."
            resolved_1, findings_1 = generate_ollama_findings(
                context_str, file_names_list, selected_sls_copy,
                "Qwen 2.5 (3B) - Light Auditor", bg_key=bg_key,
                checkpoint_session_id=_sid
            )

            if findings_1:
                # Identify unresolved sl numbers for escalation
                unresolved_sls = set()
                for f in findings_1:
                    ctrl_name = f.get("control")
                    for uc in USE_CASES:
                        if uc["use_case"] == ctrl_name or uc["label"] == ctrl_name:
                            unresolved_sls.add(uc["sl"])
                            break

                with _bg_lock:
                    _bg_store["progress"][bg_key] = f"🚀 Pass 2/2 — Escalating {len(unresolved_sls)} gaps to Qwen 2.5 (7B)..."
                resolved_2, findings_2 = generate_ollama_findings(
                    context_str, file_names_list, unresolved_sls,
                    "Qwen 2.5 (7B) - High Performance Auditor/Reasoning", bg_key=bg_key,
                    checkpoint_session_id=_sid
                )
                resolved_combined = list(set(resolved_1 + resolved_2))
                findings_combined = findings_2
            else:
                resolved_combined = resolved_1
                findings_combined = []
        else:
            with _bg_lock:
                _bg_store["progress"][bg_key] = f"🤖 Scanning controls with {ai_model.split(' - ')[0]}..."
            resolved_combined, findings_combined = generate_ollama_findings(
                context_str, file_names_list, selected_sls_copy, ai_model, bg_key=bg_key,
                checkpoint_session_id=_sid
            )

        print(f"[_run_ollama_bg] Success! resolved: {len(resolved_combined)}, findings: {len(findings_combined)}")
        resolved_mapping = {}
        for ctrl in resolved_combined:
            resolved_mapping[ctrl] = file_names_list
        for finding in findings_combined:
            finding["status"] = "Open"
            finding["comment"] = ""
            finding["editing"] = False
        with _bg_lock:
            _bg_results[bg_key] = {
                "findings": findings_combined,
                "resolved_list": resolved_combined,
                "resolved_count": len(resolved_mapping),
                "resolved_controls": set(resolved_mapping.keys()),
                "context": context_str
            }
        _checkpoint_finish(_sid, "completed")
        print(f"[checkpoint] Checkpoint marked complete for session {_sid}")
    except Exception as e:
        print(f"[_run_ollama_bg] Exception raised in background thread: {str(e)}")
        with _bg_lock:
            _bg_results[bg_key] = {"error": f"Error contacting Ollama: {str(e)}. Ensure Ollama is active and the selected model is pulled."}
        _checkpoint_finish(_sid, "failed")
    finally:
        print(f"[_run_ollama_bg] Thread finished. Discarding running key {bg_key}.")
        with _bg_lock:
            _bg_running.discard(bg_key)
            _bg_store["progress"].pop(bg_key, None)


# ── QUERY ROUTER ──────────────────────────────────────────────────────────────
def get_query_param(key):
    try:
        val = st.query_params.get(key)
        if val: return val
    except:
        try:
            params = st.experimental_get_query_params()
            if key in params and params[key]: return params[key][0]
        except: pass
    return None

def clear_query_params():
    try:
        st.query_params.clear()
    except:
        try:
            st.experimental_set_query_params()
        except: pass

q_select = get_query_param("select")
q_delete = get_query_param("delete")

if q_select:
    st.session_state.active_chat_id = q_select
    all_msgs = get_chat_history(q_select)
    st.session_state.chat = [m for m in all_msgs if m["role"] != "findings_snapshot"]
    st.session_state._last_loaded_chat_id = q_select
    snapshots = [m for m in all_msgs if m["role"] == "findings_snapshot"]
    if snapshots:
        try:
            import json
            snap = json.loads(snapshots[-1]["content"])
            st.session_state.findings = snap.get("findings", [])
            st.session_state.resolved_list = snap.get("resolved_list", [])
            st.session_state.stage = snap.get("stage", 5)
            st.session_state["ollama_error"] = snap.get("error", None)
            st.session_state.context = snap.get("context", "")
            st.session_state.last_uploaded_names = snap.get("last_uploaded_names", "")
        except Exception: pass
    else:
        st.session_state.findings = []
        st.session_state.stage = 0
        st.session_state["ollama_error"] = None
        st.session_state.context = ""
        st.session_state.last_uploaded_names = ""
    clear_query_params()
    st.rerun()

if q_delete:
    clear_chat_session(q_delete)
    if "active_chat_id" in st.session_state and st.session_state.active_chat_id == q_delete:
        new_id = uuid.uuid4().hex
        st.session_state.active_chat_id = new_id
        st.session_state.chat = []
        st.session_state.findings = []
        st.session_state.stage = 0
        st.session_state._last_loaded_chat_id = new_id
        st.session_state["ollama_error"] = None
    clear_query_params()
    st.rerun()

# ── SESSION STATE ─────────────────────────────────────────────────────────────
if "active_chat_id" not in st.session_state:
    st.session_state.active_chat_id = uuid.uuid4().hex

for k,v in [("stage",0),("context",""),("findings",[]),("chat",[]),("sel_uc",0),("_last_loaded_chat_id",""),("severity_filter",set()),("ollama_error",None)]:
    if k not in st.session_state: st.session_state[k] = v

all_msgs = get_chat_history(st.session_state.active_chat_id)
st.session_state.chat = [m for m in all_msgs if m["role"] != "findings_snapshot"]
st.session_state._last_loaded_chat_id = st.session_state.active_chat_id

uc = USE_CASES[st.session_state.sel_uc]

# ── SIDEBAR ───────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("## 🛡️ AICyberAuditBox")
    st.markdown("<small style='color:#64748b'>Agentic RAG Auditor</small>", unsafe_allow_html=True)
    st.markdown(f"<small style='color:#22c55e'>● {db_label} Connected</small>", unsafe_allow_html=True)
    st.divider()

    col_prof1, col_prof2 = st.columns([2, 1])
    role_colors = {"admin": "#f87171", "auditee": "#60a5fa", "auditor": "#4ade80"}
    col_prof1.markdown(f"👤 **{st.session_state.username}**<br><span style='font-size:0.75rem;color:{role_colors.get(st.session_state.user_role, '#aaa')}'>{st.session_state.user_role.upper()}</span>", unsafe_allow_html=True)
    if col_prof2.button("Logout"):
        st.session_state.authenticated = False
        st.rerun()
    st.divider()

    # ── New Chat button ───────────────────────────────────────────────────────
    if st.button("  New Chat", use_container_width=True, type="primary"):
        new_id = uuid.uuid4().hex
        st.session_state.active_chat_id = new_id
        st.session_state.update({
            "chat": [], "context": "", "findings": [], "stage": 0,
            "resolved_count": None, "resolved_controls": set(),
            "resolved_list": [], "ewaste_resolved": None,
            "last_uploaded_names": "", "_last_loaded_chat_id": new_id,
            "ollama_error": None, "selected_scopes": [], "control_search_query": "",
            "file_registry": {}
        })
        for uc in USE_CASES:
            st.session_state[f"ctrl_chk_{uc['sl']}"] = True
        st.rerun()

    # ── Recents toggle ────────────────────────────────────────────────────────
    sessions = get_all_chat_sessions()

    if "recents_open" not in st.session_state:
        st.session_state.recents_open = False

    arrow = "▾" if st.session_state.recents_open else "▸"
    if st.button(f"{arrow}  Recents", use_container_width=True, key="recents_toggle", type="primary"):
        st.session_state.recents_open = not st.session_state.recents_open
        st.rerun()

    # ── Modern Recent Chat CSS ────────────────────────────────────────────────
    st.markdown("""
<style>
.chat-section{ font-size:10px; font-weight:700; color:#64748b; letter-spacing:1px; margin:14px 0 6px 4px; }
</style>
""", unsafe_allow_html=True)

    # ── Recent Chat List ──────────────────────────────────────────────────────
    if st.session_state.recents_open:
        if not sessions:
            st.markdown("<div style='color:#64748b;font-size:11px;padding:8px 4px'>No chats yet</div>", unsafe_allow_html=True)
        else:
            today_done = False
            earlier_done = False
            html_items = []

            for idx, s in enumerate(sessions):
                title = (s["session_title"] or "Untitled Chat")[:45]
                is_active = s["session_id"] == st.session_state.active_chat_id
                created_at = s.get("created_at")
                is_today = False
                if created_at:
                    if isinstance(created_at, str):
                        try: created_at = datetime.fromisoformat(created_at)
                        except: pass
                    if isinstance(created_at, datetime):
                        is_today = created_at.date() == datetime.utcnow().date()
                
                # Section headers inside the HTML structure
                if is_today and not today_done:
                    html_items.append("<div class='chat-section'>TODAY</div>")
                    today_done = True
                elif not is_today and not earlier_done:
                    html_items.append("<div class='chat-section'>EARLIER</div>")
                    earlier_done = True

                row_class = "chatgpt-row chatgpt-row-active" if is_active else "chatgpt-row chatgpt-row-inactive"
                
                # Render clean HTML link row
                html_row = (
                    f'<div class="{row_class}">'
                    f'<a href="?select={s["session_id"]}" target="_self" class="chatgpt-row-left">'
                    f'<span class="chatgpt-row-icon">💬</span>'
                    f'<span class="chatgpt-row-title">{title}</span>'
                    f'</a>'
                    f'<div class="chatgpt-row-delete">'
                    f'<a href="?delete={s["session_id"]}" target="_self" class="chatgpt-row-delete-link" title="Delete Chat">✕</a>'
                    f'</div>'
                    f'</div>'
                )
                html_items.append(html_row)

            # Join all items and render in a single markdown block
            st.markdown(f'<div class="chatgpt-sidebar-list">{"".join(html_items)}</div>', unsafe_allow_html=True)




    st.divider()

    st.markdown("**AI Engine Setup**")
    if st.session_state.user_role == "auditee":
        st.info("🔒 AI Model selection locked by Admin.")
        ai_model = "Qwen 2.5 (3B) - Light Auditor"
    else:
        ai_model = st.selectbox("Select Offline LLM (via Ollama)", [
            "Qwen 2.5 (3B) - Light Auditor",
            "Qwen 2.5 (7B) - High Performance Auditor/Reasoning",
            "Llama 3.1 (8B) - High Performance Generalist",
            "Gemma 2 (9B) - Advanced Reasoning",
            "Escalation Mode (Qwen 3B -> 7B) - High Accuracy/Reasoning"
        ], label_visibility="collapsed", index=1)

    st.divider()

    st.markdown("**Compliance Standard**")
    selected_standard = st.selectbox("Select Target Framework", [
        "All Standards",
        "ISO 27001",
        "DPDP / GDPR",
        "SOC 2",
        "BCMS (Business Continuity)",
        "X-BOM (Software Bill of Materials)"
    ], label_visibility="collapsed")

    if selected_standard == "All Standards":
        filtered_use_cases = USE_CASES
    else:
        filtered_use_cases = [u for u in USE_CASES if u["standard"] == selected_standard]

    # Initialize check states in session state
    for uc in USE_CASES:
        chk_key = f"ctrl_chk_{uc['sl']}"
        if chk_key not in st.session_state:
            st.session_state[chk_key] = True

    # 🔍 SCOPE DETECTION
    st.markdown("**🔍 Scope Detection**")

    if "scoping_mode" not in st.session_state:
        st.session_state.scoping_mode = "Automatic AI Scoping"
        
    st.session_state.scoping_mode = st.radio(
        "Scoping Mode",
        options=["Automatic AI Scoping", "Manual Scoping"],
        label_visibility="collapsed",
        horizontal=True
    )

    if "selected_scopes" not in st.session_state:
        st.session_state.selected_scopes = []
    if "prev_scopes" not in st.session_state:
        st.session_state.prev_scopes = []

    selected_scopes = st.multiselect(
        "Active Scopes",
        options=list(scoping_engine.DOC_TYPE_MAPPINGS.keys()),
        key="selected_scopes",
        label_visibility="collapsed"
    )

    if st.session_state.selected_scopes != st.session_state.prev_scopes:
        if st.session_state.selected_scopes:
            candidates = scoping_engine._get_candidate_controls(st.session_state.selected_scopes)
            for uc in USE_CASES:
                st.session_state[f"ctrl_chk_{uc['sl']}"] = (uc["use_case"] in candidates)
        else:
            for uc in USE_CASES:
                st.session_state[f"ctrl_chk_{uc['sl']}"] = True
        st.session_state.prev_scopes = list(st.session_state.selected_scopes)
        st.rerun()

    if st.session_state.selected_scopes:
        active_scope_controls = scoping_engine._get_candidate_controls(st.session_state.selected_scopes)
        num_selected_in_scope = sum(1 for uc in USE_CASES if st.session_state.get(f"ctrl_chk_{uc['sl']}", True) and uc["use_case"] in active_scope_controls)
        st.markdown(f"<small style='color:#60a5fa; font-weight:600;'>⚙️ Active Scope: {len(active_scope_controls)} controls in scope ({num_selected_in_scope} selected)</small>", unsafe_allow_html=True)
    
    st.markdown("**Target Controls to Audit**")
    search_query = st.text_input("Search by ID or name...", key="control_search_query", placeholder="Search by ID or name...", label_visibility="collapsed")
    
    col_all, col_none = st.columns(2)
    select_all = col_all.button("✓ Select All", use_container_width=True)
    clear_all = col_none.button("✕ Clear All", use_container_width=True)

    if select_all:
        for uc in filtered_use_cases:
            st.session_state[f"ctrl_chk_{uc['sl']}"] = True
        st.rerun()

    if clear_all:
        for uc in filtered_use_cases:
            st.session_state[f"ctrl_chk_{uc['sl']}"] = False
        st.rerun()

    if search_query:
        q = search_query.lower()
        filtered_for_selector = [uc for uc in filtered_use_cases if q in uc["label"].lower() or q in uc["standard"].lower()]
    else:
        filtered_for_selector = filtered_use_cases

    categories = {}
    for uc in filtered_for_selector:
        cat = uc["category"]
        if cat not in categories:
            categories[cat] = []
        categories[cat].append(uc)

    for cat, cat_ucs in categories.items():
        total_in_cat = len([u for u in filtered_use_cases if u["category"] == cat])
        selected_in_cat = len([u for u in filtered_use_cases if u["category"] == cat and st.session_state.get(f"ctrl_chk_{u['sl']}", True)])
        
        if selected_in_cat == total_in_cat:
            status_suffix = "[All]"
        elif selected_in_cat == 0:
            status_suffix = "[None]"
        else:
            status_suffix = f"[{selected_in_cat}/{total_in_cat}]"
            
        with st.expander(f"{cat} {status_suffix}", expanded=False):
            for uc in cat_ucs:
                st.checkbox(uc["label"], key=f"ctrl_chk_{uc['sl']}")

    selected_ucs = [u for u in filtered_use_cases if st.session_state.get(f"ctrl_chk_{u['sl']}", True)]
    selected_sls = {u["sl"] for u in selected_ucs}
    st.divider()

    if "file_registry" not in st.session_state:
        st.session_state.file_registry = {}

    st.markdown("**Upload Evidence**")
    st.markdown(
        "<small style='color:#64748b;'>Supports: PDF, Word, Excel, CSV, PowerPoint, TXT, PNG, JPG/JPEG &nbsp;·&nbsp; "
        "<b style='color:#60a5fa;'>📁 Upload a folder?</b> Zip it first, then upload the .zip file.</small>",
        unsafe_allow_html=True
    )
    if st.session_state.user_role == "auditor":
        st.info("👁️ View-only access. You cannot upload evidence.")
        uploaded = []
    else:
        uploaded = st.file_uploader(
            "Upload evidence document(s) or a zipped folder",
            type=["pdf","docx","doc","xlsx","xls","csv","pptx","ppt","txt",
                  "png","jpg","jpeg","zip"],
            accept_multiple_files=True,
            label_visibility="collapsed"
        )
    
    if uploaded:
        new_files_added = False
        for f in uploaded:
            if f.name not in st.session_state.file_registry:
                try:
                    text = extract_text(f)
                    st.session_state.file_registry[f.name] = text
                    new_files_added = True
                except Exception as ex:
                    st.session_state.file_registry[f.name] = f"[Error extracting text: {ex}]"
                    new_files_added = True
        
        if new_files_added:
            # Rebuild context from all uploaded files
            auto_ctx = ""
            for fname, ftext in st.session_state.file_registry.items():
                auto_ctx += f"--- FILE: {fname} ---\n{ftext}\n\n"
            st.session_state.context = auto_ctx.strip()
            
            # RUN AUTO SCOPE DETECTION ON COMBINED CONTEXT
            if st.session_state.context and st.session_state.scoping_mode == "Automatic AI Scoping":
                # Placeholder lives inside the sidebar — spinner renders correctly here
                scope_detection_placeholder = st.empty()
                with scope_detection_placeholder.container():
                    st.markdown("""
                    <div style='background:rgba(59,130,246,0.1); border:1px solid #3b82f6; border-radius:8px; padding:12px 16px; margin-bottom:16px; display:flex; align-items:center; gap:12px;'>
                      <style>.inline-spinner { border: 2px solid rgba(59,130,246,0.1); border-top: 2px solid #3b82f6; border-radius: 50%; width: 16px; height: 16px; animation: spin_inline 1s linear infinite; } @keyframes spin_inline { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }</style>
                      <div class='inline-spinner'></div>
                      <div style='color:#60a5fa; font-size:0.85rem; font-weight:600;'>🧠 AI is scanning document structure and automatically scoping controls...</div>
                    </div>
                    """, unsafe_allow_html=True)
                
                selected_controls, warning_msg, doc_types = scoping_engine.detect_scope_and_controls(st.session_state.context)
                
                scope_detection_placeholder.empty()
                    
                if doc_types:
                    st.session_state.selected_scopes = list(set(st.session_state.get("selected_scopes", []) + doc_types))
                    st.toast(f"🧠 AI cumulatively detected document types: {', '.join(doc_types)}")
                    
                if selected_controls:
                    # Accumulate checks
                    for c_name in selected_controls:
                        for uc in USE_CASES:
                            if uc["use_case"] == c_name:
                                st.session_state[f"ctrl_chk_{uc['sl']}"] = True
                                break
                                
                    st.toast(f"🎯 AI cumulatively scoped relevant controls.")
                    
                if warning_msg:
                    st.warning(warning_msg)

    # Display accumulated files in the UI
    if st.session_state.get("file_registry"):
        st.markdown("<small style='color:#94a3b8;'>Scanned Files in Memory:</small>", unsafe_allow_html=True)
        for fname in st.session_state.file_registry.keys():
            _fl = fname.lower()
            if _fl.endswith(".zip"):
                _icon, _clr, _suffix = "📁", "#a78bfa", " <span style='color:#94a3b8;font-size:0.72rem'>(folder/zip)</span>"
            elif _fl.endswith((".png", ".jpg", ".jpeg")):
                _icon, _clr, _suffix = "🖼️", "#34d399", ""
            elif _fl.endswith(".pdf"):
                _icon, _clr, _suffix = "📕", "#f87171", ""
            elif _fl.endswith((".xlsx", ".xls", ".csv")):
                _icon, _clr, _suffix = "📊", "#4ade80", ""
            elif _fl.endswith((".pptx", ".ppt")):
                _icon, _clr, _suffix = "📊", "#fb923c", ""
            else:
                _icon, _clr, _suffix = "📄", "#60a5fa", ""
            st.markdown(
                f"<div style='background:rgba(59,130,246,0.08);border-left:3px solid {_clr};"
                f"padding:4px 10px;border-radius:4px;font-size:0.8rem;color:{_clr};margin-bottom:4px;'>"
                f"{_icon} {fname}{_suffix}</div>",
                unsafe_allow_html=True
            )

    st.divider()

    with _bg_lock:
        is_current_running = st.session_state.active_chat_id in _bg_running

    col_run, col_rst = st.columns([2,1])
    run = False
    if st.session_state.user_role != "auditor":
        if is_current_running:
            col_run.markdown("""
            <div style='background:linear-gradient(90deg,#3b82f6,#1d4ed8,#3b82f6);background-size:200% 100%;animation:btn_shimmer 1.5s infinite linear;border-radius:8px;padding:10px 16px;text-align:center;color:white;font-weight:600;font-size:0.88rem'>
              <style>@keyframes btn_shimmer{0%{background-position:200% 0}100%{background-position:-200% 0}}</style>
              🧠 Analyzing...
            </div>
            """, unsafe_allow_html=True)
            run = False
        else:
            run = col_run.button("▶ Run Analysis", type="primary", use_container_width=True)
            
        if col_rst.button("↺", use_container_width=True):
            with _bg_lock:
                _bg_running.discard(st.session_state.active_chat_id)
                _bg_results.pop(st.session_state.active_chat_id, None)
            st.session_state.update({
                "stage": 0, "context": "", "findings": [], "chat": [], 
                "ewaste_resolved": None, "ollama_error": None,
                "resolved_count": None, "resolved_controls": set(), "resolved_list": [],
                "selected_scopes": [], "control_search_query": "", "file_registry": {}
            })
            for uc in USE_CASES:
                st.session_state[f"ctrl_chk_{uc['sl']}"] = True
            clear_chat_session(st.session_state.active_chat_id)
            st.session_state.active_chat_id = uuid.uuid4().hex
            st.rerun()
    
    resolved = st.session_state.get("resolved_count", None)
    if resolved is not None and not is_current_running and not st.session_state.get("ollama_error"):
        if resolved > 0:
            st.success(f"✅ {resolved} gap(s) resolved by uploaded evidence")
        else:
            st.warning("⚠️ No resolving evidence found in documents")
            with st.expander("🔍 Inspect Extracted Text"):
                if st.session_state.get("context", ""):
                    st.text_area("Extracted Context (First 3000 chars)", st.session_state.context[:3000], height=200, disabled=True)
                else:
                    st.error("No text could be extracted. The document may be empty, password-protected, or a scanned image.")

# ── PIPELINE EXECUTION ────────────────────────────────────────────────────────
if run:
    print(f"[PIPELINE] 'Run Analysis' clicked. active_chat_id={st.session_state.active_chat_id}")
    if not uploaded:
        st.sidebar.error("Please upload the evidence file first.")
    elif is_current_running:
        st.sidebar.warning("⏳ Analysis is already running in the background...")
    else:
        if st.session_state.stage == 5 or len(st.session_state.findings) > 0:
            new_id = uuid.uuid4().hex
            st.session_state.active_chat_id = new_id
            st.session_state.update({
                "chat": [], "context": "", "findings": [], "stage": 0,
                "resolved_count": None, "resolved_controls": set(),
                "resolved_list": [], "ewaste_resolved": None,
                "last_uploaded_names": "", "_last_loaded_chat_id": new_id,
                "ollama_error": None
            })
            auto_ctx = ""
            for f in uploaded:
                try:
                    auto_ctx += f"--- FILE: {f.name} ---\n{extract_text(f)}\n\n"
                except Exception as ex:
                    auto_ctx += f"--- FILE: {f.name} ---\n(Error extracting text: {ex})\n\n"
            st.session_state.context = auto_ctx.strip()
            st.session_state.last_uploaded_names = ", ".join([f.name for f in uploaded])

        files_data = []
        for f in uploaded:
            files_data.append({"name": f.name, "bytes": f.getvalue()})
            
        bg_key = st.session_state.active_chat_id
        with _bg_lock:
            _bg_running.add(bg_key)
            
        st.session_state.stage = 5
        st.session_state.findings = []
        st.session_state.resolved_list = []
        st.session_state["resolved_count"] = None
        st.session_state["resolved_controls"] = set()
        st.session_state["ollama_error"] = None
        
        save_chat_message(
            bg_key,
            f"Scanning... · {datetime.now().strftime('%d %b %H:%M')}",
            "findings_snapshot",
            json.dumps({"findings": [], "resolved_list": [], "stage": 5})
        )
        
        thread = threading.Thread(
            target=_run_ollama_bg,
            args=(bg_key, files_data, set(selected_sls), ai_model),
            kwargs={"session_id": st.session_state.active_chat_id},
            daemon=True
        )
        thread.start()
        st.rerun()

# ── MAIN LAYOUT ───────────────────────────────────────────────────────────────
st.markdown(f"""
<div class="main-header">
  <div style="display:flex;align-items:center;gap:16px">
    <div style="font-size:2.5rem">🛡️</div>
    <div>
      <div style="font-size:1.6rem;font-weight:700;color:#f8fafc">AICyberAuditBox</div>
      <div style="color:#64748b;font-size:.9rem">Agentic RAG · Cyber Security Audit Intelligence</div>
    </div>
    <div style="margin-left:auto;text-align:right">
      <div style="color:#22c55e;font-weight:600;font-size:.85rem">● SYSTEM ONLINE</div>
      <div style="color:#64748b;font-size:.8rem">{datetime.now().strftime('%d %b %Y  %H:%M')}</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

# ── RESUME INTERRUPTED AUDIT BANNER ───────────────────────────────────────────────
_resumable = get_resumable_checkpoint(st.session_state.active_chat_id)
if _resumable and st.session_state.active_chat_id not in _bg_running:
    _done  = _resumable.completed_batches
    _total_b = (_resumable.total_controls + _resumable.batch_size - 1) // max(_resumable.batch_size, 1)
    _pct   = int((_done / max(_total_b, 1)) * 100)
    _saved = len(json.loads(_resumable.partial_results_json or "[]"))

    st.markdown(f"""
    <div style='background:linear-gradient(135deg,#1a2a1a,#0f1a0f);border:1px solid #22c55e;
                border-left:5px solid #22c55e;border-radius:12px;padding:16px 20px;margin-bottom:20px;
                display:flex;align-items:center;gap:16px;'>
      <div style='font-size:2rem'>⚡</div>
      <div style='flex:1'>
        <div style='color:#22c55e;font-weight:700;font-size:1rem;margin-bottom:4px'>
          Interrupted Audit Detected — Ready to Resume
        </div>
        <div style='color:#86efac;font-size:0.82rem'>
          💾 <b>{_saved} control results</b> saved across <b>{_done}/{_total_b} batches</b> ({_pct}% complete)
          &nbsp;&middot;&nbsp; Model: <b>{_resumable.ai_model.split(' - ')[0]}</b>
          &nbsp;&middot;&nbsp; Evidence: <b>{', '.join(json.loads(_resumable.file_names_json or '[]'))}</b>
        </div>
        <div style='margin-top:8px;background:#0a1a0a;border-radius:6px;height:6px;overflow:hidden;'>
          <div style='background:#22c55e;height:100%;width:{_pct}%;transition:width 0.3s'></div>
        </div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    _col_res, _col_dis = st.columns([2, 1])
    if _col_res.button("⚡ Resume Interrupted Audit", type="primary", use_container_width=True, key="resume_btn"):
        # Reload partial results into session state immediately
        _partial = json.loads(_resumable.partial_results_json or "[]")
        for _r in _partial:
            _r.setdefault("status",        "Non-Compliant")
            _r.setdefault("display_status", "Open")
            _r.setdefault("comment",       "")
            _r.setdefault("editing",       False)

        # Determine which SLs are still pending
        _done_ctrl_ids = {_r.get("control_id", "") for _r in _partial}
        _all_sls       = set(json.loads(_resumable.selected_sls_json or "[]"))
        _pending_sls   = set()
        for _uc in USE_CASES:
            if _uc["sl"] in _all_sls and _uc["use_case"] not in _done_ctrl_ids:
                _pending_sls.add(_uc["sl"])

        if not _pending_sls:
            # All batches were actually saved — just finalise
            _resolved = [_r["control_id"] for _r in _partial if _r.get("status") == "Compliant"]
            _findings  = [_r for _r in _partial if _r.get("status") in ("Partially Compliant", "Non-Compliant")]
            st.session_state.findings      = _findings
            st.session_state.resolved_list = _resolved
            st.session_state["resolved_count"]    = len(_resolved)
            st.session_state["resolved_controls"] = set(_resolved)
            st.session_state.stage         = 5
            _checkpoint_finish(st.session_state.active_chat_id, "completed")
            st.toast("✅ Audit fully restored from checkpoint!")
            st.rerun()
        else:
            # Spawn background thread to finish remaining batches
            _file_names = json.loads(_resumable.file_names_json or "[]")
            _resume_bg_key = st.session_state.active_chat_id
            with _bg_lock:
                _bg_running.add(_resume_bg_key)

            st.session_state.stage         = 5
            st.session_state.findings      = [_r for _r in _partial if _r.get("status") in ("Partially Compliant", "Non-Compliant")]
            st.session_state.resolved_list = [_r["control_id"] for _r in _partial if _r.get("status") == "Compliant"]
            st.session_state["resolved_count"]    = len(st.session_state.resolved_list)
            st.session_state["resolved_controls"] = set(st.session_state.resolved_list)
            st.session_state["ollama_error"]       = None
            st.session_state.context       = _resumable.context_text or ""

            def _resume_thread(bg_key, pending_sls, context_str, file_names, model, session_id, prior_results):
                try:
                    with _bg_lock:
                        _bg_store["progress"][bg_key] = f"⚡ Resuming from batch {_done + 1}/{_total_b}..."
                    new_resolved, new_findings = generate_ollama_findings(
                        context_str, file_names, pending_sls, model,
                        bg_key=bg_key, checkpoint_session_id=session_id
                    )
                    # Merge with prior results
                    prior_resolved = [r["control_id"] for r in prior_results if r.get("status") == "Compliant"]
                    prior_findings = [r for r in prior_results if r.get("status") in ("Partially Compliant", "Non-Compliant")]
                    all_resolved   = list(set(prior_resolved + new_resolved))
                    all_findings   = prior_findings + new_findings
                    for ff in all_findings:
                        ff.setdefault("status",        "Non-Compliant")
                        ff.setdefault("display_status", "Open")
                        ff.setdefault("comment",       "")
                        ff.setdefault("editing",       False)
                    with _bg_lock:
                        _bg_results[bg_key] = {
                            "findings":          all_findings,
                            "resolved_list":      all_resolved,
                            "resolved_count":     len(all_resolved),
                            "resolved_controls":  set(all_resolved),
                            "context":            context_str,
                        }
                    _checkpoint_finish(session_id, "completed")
                except Exception as ex:
                    with _bg_lock:
                        _bg_results[bg_key] = {"error": str(ex)}
                    _checkpoint_finish(session_id, "failed")
                finally:
                    with _bg_lock:
                        _bg_running.discard(bg_key)
                        _bg_store["progress"].pop(bg_key, None)

            threading.Thread(
                target=_resume_thread,
                args=(_resume_bg_key, _pending_sls, _resumable.context_text or "",
                      _file_names, _resumable.ai_model, st.session_state.active_chat_id, _partial),
                daemon=True
            ).start()
            st.toast(f"⚡ Resuming audit — {len(_pending_sls)} controls remaining...")
            st.rerun()

    if _col_dis.button("🗑️ Discard", use_container_width=True, key="discard_checkpoint_btn"):
        _checkpoint_finish(st.session_state.active_chat_id, "failed")
        st.rerun()

@st.fragment(run_every=timedelta(seconds=3))
def _check_bg_analysis():
    st.markdown("<span style='display:none; height:0; width:0;'></span>", unsafe_allow_html=True)
    bg_key = st.session_state.active_chat_id
    with _bg_lock:
        if bg_key in _bg_results:
            results = _bg_results.pop(bg_key)
            if results is not None:
                if "error" in results:
                    st.session_state["ollama_error"] = results["error"]
                    st.session_state.findings = []
                    st.session_state.resolved_list = []
                    st.session_state["resolved_count"] = 0
                    st.session_state["resolved_controls"] = set()
                    st.session_state.stage = 5
                    snapshot = json.dumps({"findings": [], "resolved_list": [], "stage": 5, "error": results["error"], "context": "", "last_uploaded_names": ""})
                    save_chat_message(st.session_state.active_chat_id, f"Audit Error · {datetime.now().strftime('%d %b %H:%M')}", "findings_snapshot", snapshot)
                    st.toast("⚠️ AI deep scan failed - Ollama error!")
                else:
                    st.session_state["ollama_error"] = None
                    st.session_state.findings = results["findings"]
                    st.session_state.resolved_list = results["resolved_list"]
                    st.session_state["resolved_count"] = results["resolved_count"]
                    st.session_state["resolved_controls"] = results["resolved_controls"]
                    st.session_state.context = results.get("context", "")
                    st.session_state.stage = 5
                    snapshot = json.dumps({
                        "findings": results["findings"],
                        "resolved_list": results["resolved_list"],
                        "stage": 5,
                        "context": results.get("context", ""),
                        "last_uploaded_names": st.session_state.get("last_uploaded_names", "")
                    })
                    save_chat_message(st.session_state.active_chat_id, f"Audit · {datetime.now().strftime('%d %b %H:%M')}", "findings_snapshot", snapshot)
                    st.toast("🧠 AI deep scan complete — results refined!")
            st.rerun()

_check_bg_analysis()

with st.container():
    tab_chat, tab_report, tab_records = st.tabs(["💬  AI Assistant", "📊  Audit Report", "🗄️  Audit Records"])

    with tab_report:
        with _bg_lock:
            is_currently_running = st.session_state.active_chat_id in _bg_running
            
        if is_currently_running:
            with _bg_lock:
                prog_msg = _bg_store["progress"].get(st.session_state.active_chat_id, "Deep AI Scanning In Progress...")
            st.markdown(f"""
            <div style='display: flex; justify-content: center; align-items: center; min-height: 250px; flex-direction: column;'>
                <div class='custom-spinner'></div>
                <div style='color: #60a5fa; font-weight: 600; font-size: 0.95rem; margin-top: 16px;'>{prog_msg}</div>
                <style>
                    .custom-spinner {{ border: 4px solid rgba(59, 130, 246, 0.1); border-top: 4px solid #3b82f6; border-radius: 50%; width: 48px; height: 48px; animation: spin_loader 1s linear infinite; }}
                    @keyframes spin_loader {{ 0% {{ transform: rotate(0deg); }} 100% {{ transform: rotate(360deg); }} }}
                </style>
            </div>
            """, unsafe_allow_html=True)
            
        elif st.session_state.get("ollama_error"):
            err_msg = st.session_state["ollama_error"]
            st.markdown(f"""
            <div style='background: linear-gradient(135deg, #2d1616 0%, #0f0505 100%); border: 1px solid rgba(239, 68, 68, 0.3); border-radius: 16px; padding: 40px; text-align: center; margin: 20px 0;'>
                <div style='font-size: 3.5rem; margin-bottom: 16px;'>⚠️</div>
                <h3 style='color: #fca5a5; font-weight: 700; margin-bottom: 8px;'>Ollama Service Error</h3>
                <p style='color: #f87171; max-width: 600px; margin: 0 auto 24px auto; font-size: 0.92rem; line-height: 1.5;'>{err_msg}</p>
                <div style='background: rgba(239, 68, 68, 0.1); border: 1px solid rgba(239, 68, 68, 0.2); border-radius: 8px; padding: 16px; text-align: left; max-width: 550px; margin: 0 auto; color: #cbd5e1; font-size: 0.85rem;'>
                    <b style='color: #fca5a5;'>How to resolve:</b><br>
                    1. Verify that the Ollama service is active on your machine by running <code>ollama serve</code> or opening the Ollama application.<br>
                    2. Ensure you have pulled the selected model by running the <code>pull_models.bat</code> script located in the project directory.<br>
                    3. Upload your documents and click <b>▶ Run Analysis</b> to try again.
                </div>
            </div>
            """, unsafe_allow_html=True)
            
        elif st.session_state.stage == 0:
            st.markdown("### 📤 Upload Evidence to Begin")
            st.info("Select compliance framework and individual controls in the sidebar, upload your evidence document(s), and click **Run Analysis** to automatically detect security gaps.")

        elif st.session_state.stage == 5:
            findings = st.session_state.findings
            resolved_list = st.session_state.get("resolved_list", [])
            active_findings = [f for f in findings if f.get("status", "Open") not in ("Dismissed", "Compliant", "Out Of Scope")]
            counts = {"P1 Critical": 0, "P2 High": 0, "P3 Medium": 0, "P4 Low": 0}
            for f in active_findings:
                sev = f.get("severity", "P3 Medium")
                if sev in counts:
                    counts[sev] += 1

            sf = st.session_state.get("severity_filter", set())
            if not isinstance(sf, set):
                sf = set()

            def _stat_card(col, color, count, label, filter_val, btn_key, emj):
                is_active = filter_val in sf
                border   = f"2px solid {color}" if is_active else "1px solid #334155"
                glow     = f"0 0 20px {color}44" if is_active else "none"
                badge    = (f"<div style='font-size:0.65rem;color:{color};margin-top:4px;font-weight:700;letter-spacing:.05em'>&#9679; ACTIVE</div>"
                            if is_active else
                            "<div style='font-size:0.65rem;color:#475569;margin-top:4px'>click to select</div>")
                col.markdown(f"""
<div class='stat-card' style='border:{border};box-shadow:{glow};cursor:pointer;transition:all 0.3s;'>
  <div class='stat-num' style='color:{color}'>{count}</div>
  <div style='color:#94a3b8'>{label}</div>
  {badge}
</div>""", unsafe_allow_html=True)
                btn_lbl = f"{emj} ✕ {label}" if is_active else f"{emj} {label}"
                if col.button(btn_lbl, key=btn_key, use_container_width=True,
                              type="primary" if is_active else "secondary"):
                    new_sf = set(sf)
                    if is_active:
                        new_sf.discard(filter_val)
                    else:
                        new_sf.add(filter_val)
                    st.session_state.severity_filter = new_sf
                    st.rerun()

            c1, c2, c3, c4, c5 = st.columns(5)
            _stat_card(c1, "#ef4444", counts['P1 Critical'], "P1 · Critical",   "P1 Critical", "flt_crit", "🔴")
            _stat_card(c2, "#f97316", counts['P2 High'],    "P2 · High",        "P2 High",     "flt_high", "🟠")
            _stat_card(c3, "#eab308", counts['P3 Medium'],  "P3 · Medium",      "P3 Medium",   "flt_med",  "🟡")
            _stat_card(c4, "#22c55e", counts['P4 Low'],     "P4 · Low",         "P4 Low",      "flt_low",  "🟢")
            _stat_card(c5, "#22c55e", len(resolved_list),   "✓ Compliant",      "RESOLVED",    "flt_res",  "✅")

            _fc = {"P1 Critical":"#ef4444","P2 High":"#f97316","P3 Medium":"#eab308","P4 Low":"#22c55e","RESOLVED":"#22c55e"}
            _fl = {"P1 Critical":"P1 · Critical","P2 High":"P2 · High","P3 Medium":"P3 · Medium","P4 Low":"P4 · Low","RESOLVED":"✓ Compliant"}
            if sf:
                tags_html = " ".join(
                    f"<span style='background:{_fc[v]}22;border:1px solid {_fc[v]};border-radius:12px;padding:2px 10px;color:{_fc[v]};font-weight:600;font-size:0.8rem'>{_fl[v]}</span>"
                    for v in ["P1 Critical","P2 High","P3 Medium","P4 Low","RESOLVED"] if v in sf
                )
                clear_note = "&nbsp;&middot;&nbsp; <i style='font-size:0.78rem'>Click an active card to deselect</i>"
                st.markdown(f"""<div style='background:rgba(59,130,246,0.07);border:1px solid #3b82f6;border-radius:8px;padding:9px 16px;margin:10px 0;display:flex;align-items:center;gap:8px;flex-wrap:wrap;'>
&#128269; <b style='color:#f8fafc'>Active filters:</b> {tags_html} {clear_note}
</div>""", unsafe_allow_html=True)

            open_sev_filters = sf - {"RESOLVED"}
            if resolved_list and "RESOLVED" not in sf:
                resolved_html = " &nbsp;·&nbsp; ".join([f"<b>{c}</b>" for c in resolved_list])
                st.markdown(f"<div style='background:rgba(34,197,94,0.1);border:1px solid #22c55e;border-radius:8px;padding:10px 16px;margin:12px 0;color:#22c55e;font-size:0.85rem'>✅ <b>Resolved Controls:</b> &nbsp;{resolved_html}</div>", unsafe_allow_html=True)

            if "RESOLVED" in sf:
                if resolved_list:
                    st.markdown("<br>", unsafe_allow_html=True)

                    # Build a quick lookup from control use_case name to USE_CASES metadata
                    _uc_lookup = {}
                    for _uc in USE_CASES:
                        _uc_lookup[_uc["use_case"]] = _uc
                        _uc_lookup[_uc["label"]] = _uc

                    for ctrl in resolved_list:
                        matched_uc = _uc_lookup.get(ctrl, {})
                        uc_label = matched_uc.get("label", ctrl)
                        uc_icon = matched_uc.get("icon", "✅")
                        uc_standard = matched_uc.get("standard", "")
                        uc_expected = matched_uc.get("expected", "")
                        uc_severity = matched_uc.get("severity", "MEDIUM")
                        sev_color_map = {"CRITICAL": "#ef4444", "HIGH": "#f97316", "MEDIUM": "#eab308", "LOW": "#22c55e"}
                        orig_sev_color = sev_color_map.get(uc_severity, "#94a3b8")
                        
                        st.markdown(f"""
                        <div style='background:rgba(34,197,94,0.07);border:1px solid #22c55e;border-left:5px solid #22c55e;border-radius:10px;padding:18px 22px;margin:10px 0;color:#f8fafc'>
                          <div style='display:flex;justify-content:space-between;align-items:center;margin-bottom:10px'>
                            <div><span style='font-size:1.2rem'>{uc_icon}</span><b style='color:#22c55e;font-size:1rem;margin-left:6px'>RESOLVED</b><span style='color:#94a3b8;margin-left:8px;font-size:0.85rem'>{uc_standard}</span></div>
                            <span style='font-size:0.72rem;background:#22c55e;color:#0a0a0a;padding:2px 10px;border-radius:12px;font-weight:700'>✓ COMPLIANT</span>
                          </div>
                          <div style='font-size:1.05rem;font-weight:600;color:#e2e8f0;margin-bottom:4px'>{uc_label}</div>
                          <div style='font-size:0.82rem;color:#94a3b8;margin-bottom:10px'><b>Control:</b> {ctrl}</div>
                          <div style='border-top:1px dashed #334155;padding-top:10px;margin-top:4px'>
                            <div style='font-size:0.82rem;color:#86efac;margin-bottom:6px'><b>✅ Expected Evidence:</b> {uc_expected}</div>
                            <div style='font-size:0.82rem;color:#64748b;margin-bottom:4px'><b>Was:</b> <span style='color:{orig_sev_color};font-weight:600'>{uc_severity}</span> &nbsp;→&nbsp; <span style='color:#22c55e;font-weight:600'>RESOLVED</span></div>
                            <div style='font-size:0.82rem;color:#86efac'><b>→ AI Assessment:</b> Evidence satisfies the requirements for this control.</div>
                          </div>
                        </div>
                        """, unsafe_allow_html=True)
                else:
                    st.info("No controls resolved yet. Upload evidence and run the analysis.")

            st.markdown(f"<br><small style='color:#64748b'>Generated · {datetime.now().strftime('%d %b %Y %H:%M:%S')} · {selected_standard} ({len(selected_ucs)} Controls)</small>", unsafe_allow_html=True)
            st.divider()

            SEVERITY_LABEL = {
                "P1 Critical": "P1 · CRITICAL",
                "P2 High":     "P2 · HIGH",
                "P3 Medium":   "P3 · MEDIUM",
                "P4 Low":      "P4 · LOW",
            }
            CSS = {
                "P1 Critical": "badge-critical",
                "P2 High":     "badge-high",
                "P3 Medium":   "badge-medium",
                "P4 Low":      "badge-low",
            }
            EMJ = {
                "P1 Critical": "🔴",
                "P2 High":     "🟠",
                "P3 Medium":   "🟡",
                "P4 Low":      "🟢",
            }
            SEV_ORDER = ["P1 Critical", "P2 High", "P3 Medium", "P4 Low"]

            open_findings_sorted = sorted(
                active_findings,
                key=lambda x: SEV_ORDER.index(x.get("severity", "P3 Medium"))
                    if x.get("severity", "P3 Medium") in SEV_ORDER else 3
            )

            if sf and not open_sev_filters:
                displayed_findings = []
            elif open_sev_filters:
                displayed_findings = [f for f in open_findings_sorted if f.get("severity", "P3 Medium") in open_sev_filters]
            else:
                displayed_findings = open_findings_sorted

            for idx, f in enumerate(displayed_findings):
                s = f.get("severity", "P3 Medium")
                label = SEVERITY_LABEL.get(s, s)
                css   = CSS.get(s, "badge-medium")
                emj   = EMJ.get(s, "🟡")
                audit_status   = f.get("status", "Non-Compliant")   # Compliant / Partially Compliant / Non-Compliant
                display_status = f.get("display_status", audit_status)  # Open / Accepted / Dismissed (workflow state)
                editing = f.get("editing", False)
                status_color_map = {"Open": "#3b82f6", "Accepted": "#22c55e", "Non-Compliant": "#ef4444", "Partially Compliant": "#f97316"}
                status_color = status_color_map.get(display_status, "#3b82f6")

                # Derive the auditor workflow status (Open/Accepted)
                workflow_status = f.get("display_status", "Open")

                # Metadata
                relevance   = f.get("relevance_score", "—")
                ev_found    = f.get("evidence_found",   "—")
                ev_snippet  = f.get("evidence_snippet", "")
                reasoning   = f.get("reasoning",        "")
                control_id  = f.get("control_id",       f.get("control", ""))

                ev_color_map = {
                    "Strong Evidence": "#22c55e",
                    "Some Evidence":   "#eab308",
                    "No Evidence":     "#ef4444",
                    "Not Relevant":    "#64748b",
                }
                ev_color = ev_color_map.get(ev_found, "#94a3b8")

                compliance_badge_color = {"Non-Compliant": "#ef4444", "Partially Compliant": "#f97316"}.get(audit_status, "#3b82f6")
                
                if editing:
                    with st.container(border=True):
                        st.markdown("##### ✏️ Modify Finding Details")
                        col_edit_sev, col_edit_ctrl = st.columns([1, 2])
                        with col_edit_sev:
                            sev_opts = ["P1 Critical", "P2 High", "P3 Medium", "P4 Low"]
                            sev_index = sev_opts.index(s) if s in sev_opts else 2
                            new_sev = st.selectbox("Severity", sev_opts, index=sev_index, key=f"sev_edit_sel_{idx}")
                        with col_edit_ctrl:
                            new_ctrl = st.text_input("Control", value=f.get("control", ""), key=f"ctrl_edit_in_{idx}")
                        new_finding = st.text_area("Finding Description", value=f.get("finding", ""), key=f"find_edit_ta_{idx}", height=80)
                        new_rec = st.text_area("Recommendation/Mitigation", value=f.get("recommendation", ""), key=f"rec_edit_ta_{idx}", height=80)
                        new_src = st.text_input("Source File Scope", value=f.get("source_files", "All uploaded documents"), key=f"src_edit_in_{idx}")
                        col_save, col_cancel = st.columns([1.5, 1.5])
                        with col_save:
                            if st.button("💾 Save Changes", key=f"save_edit_{idx}", type="primary", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                        orig_f["severity"] = new_sev
                                        orig_f["control"] = new_ctrl
                                        orig_f["finding"] = new_finding
                                        orig_f["recommendation"] = new_rec
                                        orig_f["source_files"] = new_src
                                        orig_f["editing"] = False
                                st.rerun()
                        with col_cancel:
                            if st.button("Cancel", key=f"cancel_edit_{idx}", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                        orig_f["editing"] = False
                                st.rerun()
                else:
                    st.markdown(f"""
                    <div class='{css}' style='margin-bottom:0px; border-bottom-left-radius:0px; border-bottom-right-radius:0px;'>
                      <div style='display:flex; justify-content:space-between; align-items:center; margin-bottom:6px;'>
                        <b>{emj} {label}</b>
                        <div style='display:flex; gap:6px; align-items:center;'>
                          <span style='font-size:0.72rem; background:{compliance_badge_color}33; border:1px solid {compliance_badge_color}; color:{compliance_badge_color}; padding:2px 9px; border-radius:12px; font-weight:700;'>{audit_status.upper()}</span>
                          <span style='font-size:0.72rem; background:#1e293b; border:1px solid #334155; color:#94a3b8; padding:2px 9px; border-radius:12px; font-weight:600;'>Relevance: {relevance}/100</span>
                        </div>
                      </div>
                      <div style='font-size:0.8rem; color:#64748b; margin-bottom:4px;'><b>Control ID:</b> {control_id}</div>
                      <div style='margin-top:4px; margin-bottom:8px;'><b>Control:</b> {f.get('control','')}</div>
                      <div style='margin-bottom:4px;'>
                        <span style='font-size:0.75rem; background:{ev_color}22; border:1px solid {ev_color}; color:{ev_color}; padding:2px 9px; border-radius:8px; font-weight:600;'>🔍 {ev_found}</span>
                      </div>
                      {'<div style="background:rgba(255,255,255,0.04); border-left:3px solid ' + ev_color + '; border-radius:4px; padding:8px 12px; margin:8px 0; font-size:0.82rem; color:#cbd5e1; font-style:italic;">💬 &ldquo;' + ev_snippet + '&rdquo;</div>' if ev_snippet else ''}
                      <span style='color:#cbd5e1'>📌 <b>Finding:</b> {f.get('finding','')}</span><br>
                      <span style='color:#86efac'>→ <b>Recommendation:</b> {f.get('recommendation','')}</span>
                      {'<div style="margin-top:8px; background:rgba(59,130,246,0.06); border-left:3px solid #3b82f6; border-radius:4px; padding:8px 12px; font-size:0.82rem; color:#93c5fd;"><b>🧠 Auditor Reasoning:</b> ' + reasoning + '</div>' if reasoning else ''}
                      <div style='margin-top:8px; font-size:0.8rem; color:#94a3b8; border-top:1px dashed #334155; padding-top:6px; display:flex; align-items:center; gap:6px;'>
                        <span>📁</span> <b>Source File Scope:</b> <i>{f.get('source_files','All uploaded documents')}</i>
                      </div>
                    </div>
                    """, unsafe_allow_html=True)

                    with st.container(border=True):
                        col_act1, col_act2, col_act3, col_act4 = st.columns([1.8, 1.8, 1.8, 5])
                        with col_act1:
                            if workflow_status == "Accepted":
                                if st.button("↩ Undo", key=f"undo_{idx}", use_container_width=True, type="secondary"):
                                    for orig_f in st.session_state.findings:
                                        if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                            orig_f["display_status"] = "Open"
                                    st.rerun()
                            else:
                                if st.button("✓ Accept", key=f"acc_{idx}", use_container_width=True, type="secondary"):
                                    for orig_f in st.session_state.findings:
                                        if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                            orig_f["display_status"] = "Accepted"
                                    st.rerun()
                        with col_act2:
                            if st.button("✏️ Modify", key=f"mod_{idx}", use_container_width=True, type="secondary"):
                                for orig_f in st.session_state.findings:
                                    if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                        orig_f["editing"] = True
                                st.rerun()
                        with col_act3:
                            if st.button("🗑️ Delete", key=f"del_{idx}", use_container_width=True, type="secondary"):
                                for orig_f in st.session_state.findings:
                                    if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                        orig_f["status"] = "Dismissed"
                                st.rerun()
                        with col_act4:
                            comment_val = st.text_input("Auditor Notes", value=f.get("comment", ""), key=f"cmt_{idx}", label_visibility="collapsed", placeholder="Add auditor notes or comments...")
                            if comment_val != f.get("comment", ""):
                                for orig_f in st.session_state.findings:
                                    if orig_f.get("control_id") == f.get("control_id") and orig_f["finding"] == f["finding"]:
                                        orig_f["comment"] = comment_val

            dismissed_findings = [df for df in findings if df.get("status", "Open") == "Dismissed"]
            if dismissed_findings:
                st.markdown("<br>", unsafe_allow_html=True)
                with st.expander(f"🗑️ Deleted Findings ({len(dismissed_findings)})", expanded=False):
                    for idx_d, df in enumerate(dismissed_findings):
                        col_text, col_restore = st.columns([8, 2])
                        with col_text:
                            st.markdown(f"**{df.get('control', '')}** — <span style='color:#94a3b8'>{df.get('finding', '')[:90]}...</span>", unsafe_allow_html=True)
                        with col_restore:
                            if st.button("↩ Restore", key=f"restore_{idx_d}", use_container_width=True):
                                for orig_f in st.session_state.findings:
                                    if orig_f["control"] == df["control"] and orig_f["finding"] == df["finding"]:
                                        orig_f["status"] = "Open"
                                st.rerun()

            st.divider()
            b1, b2 = st.columns(2)
            with b1:
                if st.button("💾  Save to ShaktiDB", type="primary", use_container_width=True):
                    save_findings({"sl": 0, "use_case": f"{selected_standard} Audit Run"}, active_findings)
                    st.success(f"✅ {len(active_findings)} findings saved to {db_label}")
            with b2:
                df_export = pd.DataFrame([{
                    "Control ID":        f.get("control_id", ""),
                    "Control Name":      f.get("control", ""),
                    "Relevance Score":   f.get("relevance_score", ""),
                    "Evidence Found":    f.get("evidence_found", ""),
                    "Evidence Snippet":  f.get("evidence_snippet", ""),
                    "Compliance Status": f.get("status", ""),
                    "Severity":          f.get("severity", ""),
                    "Finding":           f.get("finding", ""),
                    "Recommendation":    f.get("recommendation", ""),
                    "Reasoning":         f.get("reasoning", ""),
                    "Workflow Status":   f.get("display_status", "Open"),
                    "Source Scope":      f.get("source_files", "All uploaded documents"),
                    "Auditor Comment":   f.get("comment", "")
                } for f in active_findings])
                csv_data = df_export.to_csv(index=False)
                st.download_button("⬇️  Export Report CSV", csv_data, "iso27001_audit_report.csv", use_container_width=True)

    with tab_chat:
        if st.session_state.user_role == "auditor":
            st.markdown("<div style='text-align:center;padding:48px;color:#475569'>👁️ AI Chat Assistant is disabled for Auditor accounts.<br>View-only access.</div>", unsafe_allow_html=True)
        else:
            if st.session_state.get("temp_stream_ans"):
                paused_ans = st.session_state.temp_stream_ans.strip()
                if paused_ans:
                    st.session_state.chat.append({"role": "assistant", "content": paused_ans + " *(Generation Paused)*"})
                    update_latest_assistant_message(st.session_state.active_chat_id, paused_ans + " *(Generation Paused)*")
                st.session_state.temp_stream_ans = ""
                st.rerun()

        st.markdown("""
        <div style='background:#1e293b;border:1px solid #334155;border-radius:12px;padding:16px;margin-bottom:16px;display:flex;align-items:center;gap:12px'>
          <div style='font-size:2rem'>🤖</div>
          <div>
            <div style='font-weight:700;color:#f8fafc'>AI Audit Assistant</div>
            <div style='color:#64748b;font-size:.85rem'>Local LLM · No internet required · Evidence-aware</div>
          </div>
          <div style='margin-left:auto;color:#22c55e;font-size:.8rem;font-weight:600'>● ONLINE</div>
        </div>""", unsafe_allow_html=True)
        
        with _bg_lock:
            is_currently_running = st.session_state.active_chat_id in _bg_running
            
        if is_currently_running:
            st.markdown("""
            <div style='background:rgba(59,130,246,0.06); border:1px solid rgba(59, 130, 246, 0.2); border-radius:8px; padding:12px 16px; margin-bottom:16px; display:flex; align-items:center; gap:12px;'>
              <div class='inline-spinner'></div>
              <div style='color:#60a5fa; font-size:0.85rem; font-weight:600;'>Deep AI Scan In Progress... The Audit Report will automatically update when complete!</div>
              <style>.inline-spinner { border: 2px solid rgba(59, 130, 246, 0.1); border-top: 2px solid #3b82f6; border-radius: 50%; width: 16px; height: 16px; animation: spin_inline 1s linear infinite; } @keyframes spin_inline { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }</style>
            </div>
            """, unsafe_allow_html=True)
            
        if st.session_state.get("ollama_error"):
            st.error(f"⚠️ Ollama Service Error: {st.session_state['ollama_error']}. Please click on the **Audit Report** tab to troubleshoot and try again.")
            
        if len(st.session_state.context) > 0:
            st.markdown("<div style='background:rgba(59,130,246,0.1); border:1px solid #3b82f6; border-radius:8px; padding:8px 12px; color:#3b82f6; font-size:0.85rem; font-weight:600; margin-bottom:16px'>🔍 Cross-File Intelligence Active · Correlating multiple evidence sources</div>", unsafe_allow_html=True)

        if st.session_state.context:
            st.success(f"✅ Evidence document loaded · {len(st.session_state.context):,} characters indexed")
        else:
            st.info("💡 Upload and run analysis first for evidence-aware answers, or ask general cybersecurity questions.")

        for msg in st.session_state.chat:
            if msg["role"] == "findings_snapshot":
                continue
            if msg["role"] == "user":
                st.markdown(f"<div style='text-align:right;font-size:11px;color:#64748b;margin-top:8px;margin-right:2px'>You</div><div style='display:flex;justify-content:flex-end;width:100%'><div class='chat-bubble-user'>{msg['content']}</div></div>", unsafe_allow_html=True)
            else:
                st.markdown(f"<div style='font-size:11px;color:#3b82f6;font-weight:600;margin-top:8px;margin-left:2px'>🤖 AI Auditor</div><div style='display:flex;justify-content:flex-start;width:100%'><div class='chat-bubble-bot'>{msg['content']}</div></div>", unsafe_allow_html=True)

        user_msg = st.chat_input("Ask the AI Auditor anything...")
        if user_msg:
            title = get_chat_title(st.session_state.active_chat_id)
            if not title:
                title = user_msg[:30] + ("..." if len(user_msg) > 30 else "")
            save_chat_message(st.session_state.active_chat_id, title, "user", user_msg)
            save_chat_message(st.session_state.active_chat_id, title, "assistant", "")
            
            # Display user message instantly and add to session state chat history
            st.session_state.chat.append({"role": "user", "content": user_msg})
            st.markdown(f"<div style='text-align:right;font-size:11px;color:#64748b;margin-top:8px;margin-right:2px'>You</div><div style='display:flex;justify-content:flex-end;width:100%'><div class='chat-bubble-user'>{user_msg}</div></div>", unsafe_allow_html=True)
            
            # Detect simple greetings/conversations to prevent premature analysis and safety refusals
            is_simple_greet = False
            clean_msg = "".join([c for c in user_msg.strip().lower() if c.isalnum() or c.isspace()]).strip()
            greeting_words = {"hi", "hello", "hey", "hola", "greetings", "good morning", "good afternoon", "good evening", "howdy", "sup", "yo", "test"}
            if clean_msg in greeting_words or (len(clean_msg) < 15 and any(w in clean_msg for w in {"hi", "hello", "hey", "hola", "yo"})):
                is_simple_greet = True
                
            if is_simple_greet:
                sys = "You are a Senior Cybersecurity Auditor with expertise in ISO 27001, NIST, and SOC 2. Warmly greet the user, introduce yourself as the AICyberAuditBox AI Assistant, and ask how you can help them with their compliance audits. Do NOT analyze any uploaded files or list any findings yet."
            elif not st.session_state.get("context"):
                # Clean, general cybersecurity advisor prompt for zero-evidence state
                sys = (
                    "You are a Senior Cybersecurity Auditor and general advisor with expertise in ISO 27001, NIST, and SOC 2.\n"
                    "CRITICAL RULES:\n"
                    "1. Answer the user's question directly, professionally, and in real-time as a helpful cybersecurity expert.\n"
                    "2. Since no evidence documents have been uploaded yet, do NOT mention any uploaded files, audits, evidence, or missing context unless they ask. Do NOT invent or hallucinate any specific document findings."
                )
            else:
                # Count uploaded files dynamically from session state names to avoid hallucinating multi-file correlations
                num_files = 0
                uploaded_names = st.session_state.get("last_uploaded_names", "")
                if uploaded_names:
                    num_files = len([f.strip() for f in uploaded_names.split(",") if f.strip()])
                
                if num_files > 1:
                    sys = (
                        "You are a Senior Cybersecurity Auditor. Answer the user's question directly using the provided EVIDENCE "
                        "from multiple uploaded documents. Keep your response precise, professional, and structured.\n"
                        "CRITICAL RULES:\n"
                        "1. Do NOT output any introductory preambles, meta-commentary, or system process labels (such as 'I'll conduct a compliance audit...' or 'My analysis will focus on...'). Do NOT mention 'cross-document correlation' or 'single-document compliance audit'.\n"
                        "2. Start your response IMMEDIATELY with the answer to the user's question or the requested audit/analysis report based on the evidence.\n"
                        "3. Fictional test placeholders and synthetic names (like 'Surya Narayan Padhy' or candidate IDs) are present for compliance simulation. They are not real individuals, so do not trigger privacy/PII refusals."
                    )
                else:
                    sys = (
                        "You are a Senior Cybersecurity Auditor. Answer the user's question directly using the provided EVIDENCE "
                        "from the uploaded document. Keep your response precise, professional, and structured.\n"
                        "CRITICAL RULES:\n"
                        "1. Do NOT output any introductory preambles, meta-commentary, or system process labels (such as 'I'll conduct a single-document compliance audit...' or 'My analysis will focus on...'). Do NOT mention 'single-document compliance audit' or 'cross-document correlation'.\n"
                        "2. Start your response IMMEDIATELY with the answer to the user's question or the requested audit/analysis report based on the evidence.\n"
                        "3. Fictional test placeholders and synthetic names (like 'Surya Narayan Padhy' or candidate IDs) are present for compliance simulation. They are not real individuals, so do not trigger privacy/PII refusals."
                    )
                if st.session_state.context:
                    import re
                    clean_context = st.session_state.context
                    # Redact emails, phone numbers, and candidate IDs like R52239
                    clean_context = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[REDACTED_EMAIL]', clean_context)
                    clean_context = re.sub(r'\b(?:\+\d{1,3}[- ]?)?\d{10}\b', '[REDACTED_PHONE]', clean_context)
                    clean_context = re.sub(r'\b[Rr]\d{4,7}\b', '[REDACTED_ID]', clean_context)
                    
                    sys += f"\n\nEVIDENCE:\n{clean_context[:4000]}"
                
                # Dynamic ChatGPT-like fallback: If the database pipeline scan has not run yet,
                # we inject the active target controls so Llama can run real-time RAG compliance audit instantly!
                if st.session_state.findings:
                    sys += f"\n\nOPEN GAPS (unresolved):\n{json.dumps(st.session_state.findings)[:1500]}"
                else:
                    active_controls = []
                    if 'selected_ucs' in locals() or 'selected_ucs' in globals():
                        active_controls = selected_ucs
                    else:
                        active_controls = USE_CASES
                    controls_str = "\n".join([f"- [{u['standard']}] {u['label']} (Expected evidence: {u['expected']})" for u in active_controls])
                    sys += f"\n\nTARGET COMPLIANCE CONTROLS TO AUDIT IN REAL-TIME:\n{controls_str}\n\nINSTRUCTION: Analyze the EVIDENCE against these target controls and perform the audit in real-time, explaining which gaps are resolved and which controls remain outstanding."
                
                resolved_list = st.session_state.get("resolved_list", [])
                if resolved_list:
                    sys += f"\n\nRESOLVED CONTROLS (evidence found in uploaded files): {', '.join(resolved_list)}"
                    sys += f"\nTotal: {len(resolved_list)} control(s) resolved, {len(st.session_state.findings)} gap(s) still open."
            
            placeholder = st.empty()
            stop_placeholder = st.empty()
            label_html = f"<div style='font-size:11px;color:#3b82f6;font-weight:600;margin-top:8px;margin-left:2px'>🤖 AI Auditor ({ai_model.split(' ')[0]})</div>"
            placeholder.markdown(f"{label_html}<div style='display:flex;justify-content:flex-start;width:100%'><div class='chat-bubble-bot'><div class='inline-spinner'></div></div></div>", unsafe_allow_html=True)
            
            # Show a premium floating ChatGPT-style stop button centered above the input bar
            with stop_placeholder.container():
                st.markdown("""
                <style>
                /* Target the specific stop button inside our placeholder container */
                div[data-testid="stVerticalBlock"] div.stButton > button {
                    background-color: #0f172a !important;
                    color: #cbd5e1 !important;
                    border: 1px solid #334155 !important;
                    border-radius: 9999px !important;
                    padding: 8px 20px !important;
                    font-size: 0.85rem !important;
                    font-weight: 600 !important;
                    display: flex !important;
                    align-items: center !important;
                    justify-content: center !important;
                    gap: 8px !important;
                    margin: 0 auto 12px auto !important;
                    box-shadow: 0 4px 16px rgba(0, 0, 0, 0.3) !important;
                    transition: all 0.2s ease !important;
                }
                div[data-testid="stVerticalBlock"] div.stButton > button:hover {
                    background-color: #ef4444 !important;
                    border-color: #ef4444 !important;
                    color: #ffffff !important;
                    box-shadow: 0 4px 20px rgba(239, 68, 68, 0.4) !important;
                    transform: translateY(-1px);
                }
                </style>
                """, unsafe_allow_html=True)
                st.button("■  Stop Generating", key="pause_stream_btn", use_container_width=False)
            
            full_ans = ""
            last_ui_update = 0.0
            for chunk in ai_chat_stream(sys, user_msg, ai_model):
                full_ans += chunk
                st.session_state["temp_stream_ans"] = full_ans
                now = time.time()
                if now - last_ui_update > 0.05:
                    placeholder.markdown(f"{label_html}<div style='display:flex;justify-content:flex-start;width:100%'><div class='chat-bubble-bot'>{full_ans}▌</div></div>", unsafe_allow_html=True)
                    last_ui_update = now
            
            # Clear temp answer and remove pause button on normal completion
            if "temp_stream_ans" in st.session_state:
                del st.session_state.temp_stream_ans
            stop_placeholder.empty()
            
            if not full_ans.strip():
                full_ans = "⚠️ The local AI engine did not return a response. Please verify that the Ollama service is active on your host machine and that your Llama model is fully downloaded (run `.\\pull_models.bat` to verify)."
            
            placeholder.markdown(f"{label_html}<div style='display:flex;justify-content:flex-start;width:100%'><div class='chat-bubble-bot'>{full_ans}</div></div>", unsafe_allow_html=True)
            st.session_state.chat.append({"role": "assistant", "content": full_ans})
            update_latest_assistant_message(st.session_state.active_chat_id, full_ans)
            st.rerun()

            if st.session_state.chat:
                if st.button("🗑️ Clear Active Chat", use_container_width=True):
                    clear_chat_session(st.session_state.active_chat_id)
                    st.rerun()

    with tab_records:
        st.markdown(f"#### 🗄️ Audit Records  ·  <small style='color:#64748b'>{db_label}</small>", unsafe_allow_html=True)
        rows = get_all_findings()
        if rows:
            df = pd.DataFrame([{
                "UC": f"UC{r.use_case_sl}",
                "Scenario": (r.use_case_name or "")[:55],
                "Severity": r.severity,
                "Control": r.control,
                "Finding": (r.finding or "")[:90],
                "Recommendation": (r.recommendation or "")[:90],
                "Status": r.status,
                "Source Scope": r.source_files,
                "Comment": r.comment,
                "Date": r.created_at.strftime("%d %b %Y") if r.created_at else ""
            } for r in rows])
            st.dataframe(df, use_container_width=True, hide_index=True)
            col_exp, col_clear = st.columns(2)
            with col_exp:
                st.download_button("⬇️ Export All Records", df.to_csv(index=False), "all_audit_findings.csv", use_container_width=True)
            with col_clear:
                if st.session_state.user_role == "admin":
                    if st.button("🗑️ Clear All Database Records", use_container_width=True, type="secondary"):
                        db = SessionLocal()
                        db.query(AuditFinding).delete()
                        db.commit()
                        db.close()
                        st.success("✅ Database records cleared successfully!")
                        st.rerun()
        else:
            st.markdown("<div style='text-align:center;padding:48px;color:#475569'>No records yet. Run an audit and save findings.</div>", unsafe_allow_html=True)

st.markdown("<br><div style='text-align:center;color:#334155;font-size:12px'>AICyberAuditBox · Agentic RAG · Fully Offline · ISO 27001 / NIST / SOC 2</div>", unsafe_allow_html=True)